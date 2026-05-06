from fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from lxml import etree
import os
import re
from pathlib import Path

TEMPLATE_PATH = os.environ.get("SLIDE_TEMPLATE", "template.pptx")

mcp = FastMCP(
    name="slide-server",
    version="1.0.0",
    description="MCP Server per generare presentazioni PowerPoint con animazioni a partire da un template .pptx",
)


# ─────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────

def _clone_slide(prs: Presentation, layout_index: int = 0):
    """Aggiunge una slide usando il layout specificato."""
    layout_index = min(layout_index, len(prs.slide_layouts) - 1)
    layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(layout)


def _set_placeholder(slide, ph_idx: int, text: str, font_size: int | None = None):
    """Imposta il testo di un placeholder per indice idx."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            tf = ph.text_frame
            tf.clear()
            tf.text = text
            if font_size:
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(font_size)
            return


def _add_speaker_notes(slide, notes_text: str):
    """Scrive il testo nelle note relatore."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def _add_slide_transition(slide, dur_ms: int = 700):
    """
    Aggiunge una transizione Fade tra slide e slide tramite <p:transition>.

    Il tag <p:transition> viene inserito come figlio diretto di <p:sld>,
    prima di <p:timing> se presente, altrimenti in fondo.
    dur  = durata del fade in millisecondi (default 700 ms).
    """
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"

    sld_el = slide._element

    # Rimuovi eventuali transizioni già presenti
    for old in sld_el.findall(qn("p:transition")):
        sld_el.remove(old)

    transition = etree.Element(f"{{{P}}}transition", {
        "dur": str(dur_ms),
    })
    etree.SubElement(transition, f"{{{P}}}fade")

    timing_el = sld_el.find(qn("p:timing"))
    if timing_el is not None:
        sld_el.insert(list(sld_el).index(timing_el), transition)
    else:
        sld_el.append(transition)


# Contatore globale per garantire ID univoci nell'intero documento
_timing_id_counter = [1]


def _next_id() -> int:
    val = _timing_id_counter[0]
    _timing_id_counter[0] += 1
    return val


def _add_fade_animation(slide):
    """
    Aggiunge animazioni Fade-In (On Click, una per shape) usando la struttura
    OOXML canonica generata da PowerPoint 2019/365.

    FIX scomparsa animazioni:
    - Rimosso il blocco <p:set visibility=hidden> iniziale: PowerPoint lo
      interpretava come stato persistente e al rientro sulla slide ripristinava
      tutte le shape a 'hidden', causando la scomparsa degli elementi.
    - Ogni <p:animEffect> ha ora fill='hold' su tutti i livelli cTn, così
      lo stato finale (visibile) viene mantenuto dopo la fine dell'animazione.
    - Aggiunto advTm='0' sul cTn dell'animEffect per evitare avanzamento
      automatico che conflittava con il fill=hold.
    - Il <p:bldP> nel bldLst usa build='allAtOnce' per coerenza con il
      comportamento atteso da PowerPoint quando non c'è autoplay.
    """
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"

    sp_ids = [shape.shape_id for shape in slide.shapes]
    if not sp_ids:
        return

    root_id = _next_id()
    seq_id  = _next_id()

    timing = etree.Element(f"{{{P}}}timing")
    tnLst  = etree.SubElement(timing, f"{{{P}}}tnLst")

    root_par = etree.SubElement(tnLst, f"{{{P}}}par")
    root_cTn = etree.SubElement(root_par, f"{{{P}}}cTn", {
        "id":       str(root_id),
        "dur":      "indefinite",
        "fill":     "hold",
        "nodeType": "tmRoot",
        "restart":  "never",
    })
    root_child = etree.SubElement(root_cTn, f"{{{P}}}childTnLst")

    # --- Sequenza principale: un clickEffect (fade) per shape
    seq = etree.SubElement(root_child, f"{{{P}}}seq", {"concurrent": "1", "nextAc": "seek"})
    seq_cTn = etree.SubElement(seq, f"{{{P}}}cTn", {
        "id":       str(seq_id),
        "dur":      "indefinite",
        "fill":     "hold",
        "nodeType": "mainSeq",
        "restart":  "never",
    })
    seq_stCond = etree.SubElement(
        etree.SubElement(seq_cTn, f"{{{P}}}stCondLst"), f"{{{P}}}cond"
    )
    seq_stCond.set("delay", "indefinite")
    seq_child = etree.SubElement(seq_cTn, f"{{{P}}}childTnLst")

    for grp_idx, spid in enumerate(sp_ids):
        outer_id = _next_id()
        par_id   = _next_id()
        inner_id = _next_id()
        anim_id  = _next_id()

        outer_par = etree.SubElement(seq_child, f"{{{P}}}par")
        outer_cTn = etree.SubElement(outer_par, f"{{{P}}}cTn", {
            "id":      str(outer_id),
            "fill":    "hold",
            "restart": "never",
        })
        outer_stCond = etree.SubElement(
            etree.SubElement(outer_cTn, f"{{{P}}}stCondLst"), f"{{{P}}}cond"
        )
        outer_stCond.set("evt",   "onClick")
        outer_stCond.set("delay", "0")
        outer_child = etree.SubElement(outer_cTn, f"{{{P}}}childTnLst")

        click_par = etree.SubElement(outer_child, f"{{{P}}}par")
        cTn_click = etree.SubElement(click_par, f"{{{P}}}cTn", {
            "id":            str(par_id),
            "presetID":      "10",
            "presetClass":   "entr",
            "presetSubtype": "0",
            "fill":          "hold",
            "grpId":         str(grp_idx),
            "nodeType":      "clickEffect",
            "restart":       "never",
        })
        stCond = etree.SubElement(etree.SubElement(cTn_click, f"{{{P}}}stCondLst"), f"{{{P}}}cond")
        stCond.set("delay", "0")

        childTn   = etree.SubElement(cTn_click, f"{{{P}}}childTnLst")
        inner_par = etree.SubElement(childTn, f"{{{P}}}par")
        cTn_inner = etree.SubElement(inner_par, f"{{{P}}}cTn", {
            "id":      str(inner_id),
            "dur":     "500",
            "fill":    "hold",
            "restart": "never",
        })
        inner_child = etree.SubElement(cTn_inner, f"{{{P}}}childTnLst")

        anim  = etree.SubElement(inner_child, f"{{{P}}}animEffect", {
            "transition": "in",
            "filter":     "fade",
        })
        cBhvr = etree.SubElement(anim, f"{{{P}}}cBhvr")
        # fill='hold' + advTm='0': mantiene lo stato finale (visibile)
        # e impedisce l'avanzamento automatico che resettava la shape
        etree.SubElement(cBhvr, f"{{{P}}}cTn", {
            "id":    str(anim_id),
            "dur":   "500",
            "fill":  "hold",
            "advTm": "0",
        })
        tgtEl = etree.SubElement(cBhvr, f"{{{P}}}tgtEl")
        etree.SubElement(tgtEl, f"{{{P}}}spTgt", {"spid": str(spid)})

    # nextCondLst obbligatorio sulla seq
    next_cond = etree.SubElement(
        etree.SubElement(seq, f"{{{P}}}nextCondLst"), f"{{{P}}}cond"
    )
    next_cond.set("evt",   "onNext")
    next_cond.set("delay", "0")

    # --- bldLst: registra ogni shape nel build list con build='allAtOnce'
    bldLst = etree.SubElement(timing, f"{{{P}}}bldLst")
    for grp_idx, spid in enumerate(sp_ids):
        etree.SubElement(bldLst, f"{{{P}}}bldP", {
            "spid":     str(spid),
            "grpId":    str(grp_idx),
            "uiExpand": "1",
            "animBg":   "1",
            "build":    "allAtOnce",
        })

    # Sostituisci il timing esistente
    sld_el = slide._element
    for old in sld_el.findall(qn("p:timing")):
        sld_el.remove(old)
    sld_el.append(timing)


def _remove_template_slides(prs: Presentation) -> None:
    """
    Rimuove tutte le slide presenti nel template in modo sicuro,
    senza lasciare relazioni orfane nel pacchetto OOXML.
    """
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in reversed(list(sld_id_lst)):
        rId = sld_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sld_id_lst.remove(sld_id)


def _topic_to_filename(topic: str) -> str:
    """
    Converte il topic in un nome file .pptx sicuro e leggibile.
    Es: 'Intelligenza Artificiale nel 2025' -> 'Intelligenza Artificiale nel 2025.pptx'
    Rimuove caratteri non ammessi nei path di file system e tronca a 50 caratteri.
    """
    safe = re.sub(r'[\\/:*?"<>|]', "", topic).strip()
    safe = re.sub(r"\s+", " ", safe)
    safe = safe[:50].rstrip()
    return f"{safe}.pptx" if safe else "presentazione.pptx"


def _build_presentation(topic: str, slides_content: list[dict]) -> Presentation:
    """
    Costruisce una Presentation partendo dal template.
    Le note relatore sono scritte in prima persona, pronte da leggere a voce alta.
    """
    _timing_id_counter[0] = 1

    prs = Presentation(TEMPLATE_PATH)
    _remove_template_slides(prs)

    # ─── Slide 1: Titolo
    slide_title = _clone_slide(prs, 0)
    _set_placeholder(slide_title, 0, topic, font_size=40)
    _add_speaker_notes(
        slide_title,
        f"Benvenuti a questa presentazione su '{topic}'. "
        "Mi chiamo [Nome] e oggi vi guiderò attraverso i principali aspetti di questo argomento. "
        "L'obiettivo è fornirvi una panoramica chiara e operativa, con esempi concreti. "
        "Iniziamo subito.",
    )
    _add_fade_animation(slide_title)
    _add_slide_transition(slide_title)

    # ─── Slide 2: Agenda
    slide_agenda = _clone_slide(prs, min(1, len(prs.slide_layouts) - 1))
    _set_placeholder(slide_agenda, 0, "Agenda")
    agenda_items = [f"{i + 1}. {s['title']}" for i, s in enumerate(slides_content)]
    _set_placeholder(slide_agenda, 1, "\n".join(agenda_items), font_size=18)
    topic_list = ", ".join([s["title"] for s in slides_content])
    _add_speaker_notes(
        slide_agenda,
        f"Nel corso di questa presentazione affronteremo i seguenti argomenti: {topic_list}. "
        "Ogni sezione è pensata per essere autonoma, quindi se avete domande su un punto specifico "
        "potete interrompermi durante quella sezione. Partiamo dal primo argomento.",
    )
    _add_fade_animation(slide_agenda)
    _add_slide_transition(slide_agenda)

    # ─── Slide 3..N+2: Contenuto
    for sc in slides_content:
        layout_idx = min(1, len(prs.slide_layouts) - 1)
        slide = _clone_slide(prs, layout_idx)
        _set_placeholder(slide, 0, sc["title"])
        bullets = sc.get("bullets", [])
        body_text = "\n".join(f"\u2022 {b}" for b in bullets)
        _set_placeholder(slide, 1, body_text, font_size=18)
        _add_speaker_notes(slide, sc.get("notes", ""))
        _add_fade_animation(slide)
        _add_slide_transition(slide)

    # ─── Ultima Slide: Ringraziamenti
    slide_thanks = _clone_slide(prs, 0)
    _set_placeholder(slide_thanks, 0, "Grazie!", font_size=48)
    _set_placeholder(
        slide_thanks,
        1,
        "Domande? Siamo felici di rispondere.\n\nContatti: [email] | [LinkedIn]",
        font_size=20,
    )
    _add_speaker_notes(
        slide_thanks,
        f"Abbiamo concluso la presentazione su '{topic}'. "
        "Spero che i contenuti siano stati chiari e utili. "
        "Sono ora a disposizione per rispondere a qualsiasi domanda abbiate. "
        "Potete contattarmi anche dopo via email o LinkedIn ai recapiti che vedete sullo schermo.",
    )
    _add_fade_animation(slide_thanks)
    _add_slide_transition(slide_thanks)

    return prs


# ─────────────────────────────────────────────────────────────
# Tool: genera presentazione
# ─────────────────────────────────────────────────────────────

@mcp.tool(
    description="""
Genera una presentazione PowerPoint (.pptx) a partire da un argomento e da N slide di contenuto.
La presentazione avrà sempre N+3 slide:
  - Slide 1: Titolo con il nome dell'argomento
  - Slide 2: Agenda numerata con tutti i titoli
  - Slide 3..N+2: Slide di contenuto richieste
  - Slide N+3: Ringraziamenti finali

Il nome del file di output viene derivato automaticamente dal topic (es. topic='Cloud Computing' → 'Cloud Computing.pptx').
Se vuoi un percorso specifico, passa output_path esplicitamente.

Ogni slide include:
  - Animazioni Fade-In (On Click) stabili: le shape rimangono visibili dopo l'animazione
  - Transizione Fade tra slide e slide
  - Note relatore scritte in prima persona, pronte da leggere a voce alta

Convenzioni rispettate:
  - Titolo: massimo 7 parole
  - Bullet point per slide: massimo 5
  - Ogni bullet: massimo 15 parole
  - Note: testo diretto, in prima persona, senza istruzioni meta ('dì questo', 'spiega quello')

Few-shot examples
-----------------
Esempio 1 – singola slide:
  topic = "Intelligenza Artificiale nel 2025"
  slides = [
    {
      "title": "Modelli Linguistici",
      "bullets": ["GPT-4 e successori", "RAG e knowledge graph", "Costi in calo del 60%"],
      "notes": "Per quanto riguarda i modelli linguistici, negli ultimi due anni abbiamo assistito a una crescita esponenziale delle capacità di GPT-4 e dei suoi successori. Il costo per token è sceso del 60% rispetto al 2023, rendendo questi strumenti accessibili anche alle PMI. Le architetture RAG permettono oggi di collegare i modelli a knowledge graph aziendali, riducendo le allucinazioni in modo significativo."
    }
  ]

Esempio 2 – tre slide:
  topic = "Sostenibilità Aziendale"
  slides = [
    {
      "title": "Perché la Sostenibilità",
      "bullets": ["Riduzione emissioni CO2", "Risparmio energetico", "Reputazione del brand"],
      "notes": "La sostenibilità non è più un'opzione strategica, ma una necessità competitiva. Secondo il rapporto IPCC 2024, le aziende che hanno adottato piani di riduzione delle emissioni hanno registrato un risparmio energetico medio del 18% e un miglioramento della brand reputation del 34% presso i consumatori under 35."
    },
    {
      "title": "Strumenti e Standard",
      "bullets": ["ESG reporting", "ISO 14001", "Carbon footprint calculator"],
      "notes": "Gli standard ESG forniscono un framework comune per misurare e comunicare le performance ambientali. La certificazione ISO 14001 è oggi richiesta da oltre il 60% delle gare d'appalto pubbliche in Europa. I carbon footprint calculator integrati nei gestionali aziendali permettono di tracciare le emissioni in tempo reale, agevolando il reporting obbligatorio previsto dalla CSRD."
    },
    {
      "title": "Piano d'Azione",
      "bullets": ["Quick wins nei prossimi 3 mesi", "Obiettivi annuali", "KPI di monitoraggio"],
      "notes": "Nei prossimi tre mesi è possibile ottenere risultati concreti intervenendo su illuminazione, gestione dei rifiuti e ottimizzazione dei consumi energetici negli uffici. Gli obiettivi annuali devono essere SMART e collegati a KPI misurabili, come tonnellate di CO2 evitate o percentuale di energia rinnovabile utilizzata. Monitorare questi indicatori trimestralmente consente di correggere la rotta in tempo."
    }
  ]
"""
)
def generate_presentation(
    topic: str,
    slides: list[dict],
    output_path: str = "",
) -> str:
    """
    Genera una presentazione PowerPoint.

    Args:
        topic: Il titolo/argomento principale della presentazione.
        slides: Lista di dict, ognuno con:
                  - title (str): titolo della slide (max 7 parole)
                  - bullets (list[str]): elenco puntato, max 5 voci da max 15 parole
                  - notes (str): testo in prima persona da leggere come note relatore,
                                 senza istruzioni meta. Min 10 parole.
        output_path: Percorso del file .pptx da generare.
                     Se vuoto (default), viene derivato automaticamente dal topic.

    Returns:
        Messaggio con il percorso del file generato.
    """
    errors = []
    for i, s in enumerate(slides):
        title_words = len(s.get("title", "").split())
        if title_words > 7:
            errors.append(f"Slide {i+1}: il titolo ha {title_words} parole (max 7).")
        bullets = s.get("bullets", [])
        if len(bullets) > 5:
            errors.append(f"Slide {i+1}: troppi bullet ({len(bullets)}, max 5).")
        for j, b in enumerate(bullets):
            bw = len(b.split())
            if bw > 15:
                errors.append(f"Slide {i+1}, bullet {j+1}: {bw} parole (max 15).")
    if errors:
        return "Errori di validazione:\n" + "\n".join(errors)

    # Determina il percorso di output: usa output_path se fornito, altrimenti topic
    if output_path:
        out = Path(output_path)
    else:
        out = Path(_topic_to_filename(topic))

    prs = _build_presentation(topic, slides)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    total = len(slides) + 3
    return (
        f"Presentazione generata con successo: {out.resolve()}\n"
        f"File: {out.name}\n"
        f"Slide totali: {total} ({len(slides)} di contenuto + titolo + agenda + ringraziamenti)"
    )


# ─────────────────────────────────────────────────────────────
# Tool: info template
# ─────────────────────────────────────────────────────────────

@mcp.tool(
    description="""
Restituisce informazioni sul template PowerPoint attualmente configurato.
Utile per verificare quanti layout sono disponibili prima di generare una presentazione.

Few-shot examples
-----------------
Esempio 1:
  Input: nessun parametro richiesto
  Output: "Template: template.pptx | Layout disponibili: 11 | Slide già presenti: 0"

Esempio 2 (template non trovato):
  Output: "Errore: template.pptx non trovato. Assicurati che il file esista nella directory corrente."
"""
)
def template_info() -> str:
    """
    Ritorna informazioni sul template .pptx configurato.

    Returns:
        Stringa descrittiva con path, numero di layout e slide presenti.
    """
    path = Path(TEMPLATE_PATH)
    if not path.exists():
        return f"Errore: {TEMPLATE_PATH} non trovato. Assicurati che il file esista nella directory corrente."
    prs = Presentation(str(path))
    return (
        f"Template: {path.resolve()}\n"
        f"Layout disponibili: {len(prs.slide_layouts)}\n"
        f"Slide già presenti nel template: {len(prs.slides)}\n"
        f"Dimensioni slide: {prs.slide_width.inches:.1f}\" x {prs.slide_height.inches:.1f}\""
    )


# ─────────────────────────────────────────────────────────────
# Tool: valida contenuto slide
# ─────────────────────────────────────────────────────────────

@mcp.tool(
    description="""
Valida il contenuto di una lista di slide rispetto alle convenzioni standard delle presentazioni:
  - Titolo slide: max 7 parole
  - Bullet per slide: max 5
  - Parole per bullet: max 15
  - Note relatore: obbligatorie (almeno 10 parole), in prima persona

Ritorna OK se tutto è valido, oppure un elenco dettagliato degli errori.

Few-shot examples
-----------------
Esempio 1 – tutto valido:
  slides = [
    {
      "title": "Cloud Computing",
      "bullets": ["Scalabilità on-demand", "Riduzione costi IT"],
      "notes": "Il cloud computing ha trasformato radicalmente le infrastrutture IT negli ultimi dieci anni, spostando la spesa da CAPEX a OPEX e consentendo alle aziende di scalare in minuti anziché in mesi."
    }
  ]
  Output: "✅ Tutte le slide rispettano le convenzioni."

Esempio 2 – errori rilevati:
  slides = [
    {
      "title": "Questa è una slide con un titolo molto lungo che supera il limite",
      "bullets": ["b1","b2","b3","b4","b5","b6"],
      "notes": "ok"
    }
  ]
  Output:
    "❌ Errori trovati:\\n"
    "- Slide 1: titolo ha 12 parole (max 7)\\n"
    "- Slide 1: 6 bullet points (max 5)\\n"
    "- Slide 1: note troppo brevi (1 parola, min 10)"
"""
)
def validate_slides(slides: list[dict]) -> str:
    """
    Valida una lista di slide rispetto alle convenzioni.

    Args:
        slides: Lista di dict con title, bullets, notes.

    Returns:
        Stringa con esito validazione.
    """
    errors = []
    for i, s in enumerate(slides, 1):
        title = s.get("title", "")
        bullets = s.get("bullets", [])
        notes = s.get("notes", "")

        tw = len(title.split())
        if tw > 7:
            errors.append(f"Slide {i}: titolo ha {tw} parole (max 7)")

        if len(bullets) > 5:
            errors.append(f"Slide {i}: {len(bullets)} bullet points (max 5)")

        for j, b in enumerate(bullets, 1):
            bw = len(b.split())
            if bw > 15:
                errors.append(f"Slide {i}, bullet {j}: {bw} parole (max 15)")

        nw = len(notes.split())
        if nw < 10:
            errors.append(f"Slide {i}: note troppo brevi ({nw} parole, min 10)")

    if errors:
        return "❌ Errori trovati:\n" + "\n".join(f"- {e}" for e in errors)
    return "✅ Tutte le slide rispettano le convenzioni."


if __name__ == "__main__":
    mcp.run()
